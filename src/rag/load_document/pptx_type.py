from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import List, Tuple

# PPTX
from pptx import Presentation

from src.rag.load_document.base_document import BaseDocumentExtractor


class PPTXExtractor(BaseDocumentExtractor):
    def _extract_blocks(self, file_path: Path) -> List[Tuple[int, str, bytes | str]]:
        prs = Presentation(str(file_path))
        all_blocks: List[Tuple[int, str, bytes | str]] = []

        def _process_slide(
                slide_num: int, slide
        ) -> List[Tuple[int, str, bytes | str]]:
            slide_blocks: List[Tuple[int, str, bytes | str]] = []
            base_index = slide_num * 100_000
            shape_order = 0

            for shape in slide.shapes:
                idx = base_index + shape_order

                # Text frames
                if shape.has_text_frame:
                    lines = []
                    for para in shape.text_frame.paragraphs:
                        para_text = "".join(run.text for run in para.runs if run.text)
                        if para_text.strip():
                            lines.append(para_text)
                    if lines:
                        slide_blocks.append((idx, "text", "\n".join(lines)))
                        shape_order += 1

                # Tables
                if shape.has_table:
                    table_lines = []
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        table_lines.append(row_text)
                    table_str = "\n".join(table_lines)
                    if table_str.strip():
                        slide_blocks.append((idx + 1, "text", table_str))
                        shape_order += 2

                # Pictures / images
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        img_bytes = shape.image.blob
                        slide_blocks.append((idx + 2, "image", img_bytes))
                        shape_order += 3
                    except Exception as exc:
                        log.warning(
                            "PPTX slide %d shape image error: %s", slide_num, exc
                        )

                # Grouped shapes – recurse one level
                if shape.shape_type == 6:  # GROUP
                    try:
                        for sub_shape in shape.shapes:
                            if sub_shape.shape_type == 13:
                                img_bytes = sub_shape.image.blob
                                slide_blocks.append(
                                    (idx + shape_order, "image", img_bytes)
                                )
                                shape_order += 1
                    except Exception:
                        pass

            return slide_blocks

        with ThreadPoolExecutor(max_workers=self.max_workers) as pool:
            futures = {
                pool.submit(_process_slide, i, slide): i
                for i, slide in enumerate(prs.slides)
            }
            for future in as_completed(futures):
                try:
                    all_blocks.extend(future.result())
                except Exception as exc:
                    log.warning("PPTX slide processing error: %s", exc)

        return all_blocks
